from pptx import Presentation

courseName = 'D3 The Big Picture'
courseSlug = 'd3-big-picture'

data = {

    "slides": [
        {"module_num":1, "slide_title":"What is D3", "slide_layout":0},
        {"module_num":1, "slide_title":"D3 History", "slide_layout":2},
        {"module_num":1, "slide_title":"D3 Definition", "slide_layout":2},
        {"module_num":1, "slide_title":"D3 Examples", "slide_layout":2},
        {"module_num":1, "slide_title":"Where to Find More", "slide_layout":2},
        {"module_num":2, "slide_title":"Why use D3?", "slide_layout":0},
        {"module_num":2, "slide_title":"D3 Compatibility", "slide_layout":2},
        {"module_num":2, "slide_title":"D3 Performance", "slide_layout":2},
        {"module_num":2, "slide_title":"Where to Find More", "slide_layout":2},
        {"module_num":3, "slide_title":"When to Use D3", "slide_layout":0},
        {"module_num":3, "slide_title":"Web App Example", "slide_layout":2},
        {"module_num":3, "slide_title":"Data Journalism Example", "slide_layout":2},
        {"module_num":3, "slide_title":"Data Scientist Example", "slide_layout":2},
        {"module_num":3, "slide_title":"Where to Find More", "slide_layout":2},
        {"module_num":4, "slide_title":"D3 Architecture", "slide_layout":0},
        {"module_num":4, "slide_title":"D3 Library Components", "slide_layout":2},
        {"module_num":4, "slide_title":"SVG Elements", "slide_layout":2},
        {"module_num":4, "slide_title":"Effect on HTML", "slide_layout":2},
        {"module_num":4, "slide_title":"Where to Find More", "slide_layout":2},
    ]
# nested json
#     "modules": [
#         {"title": "module 1",
#          "number": "1",
#          "slides": [
#              {"title": "Slide 1", "type": "2"},
#              {"title": "Slide 2", "type": "2"},
#              {"title": "Slide 3", "type": "2"},
#              {"title": "Slide 4", "type": "2"},
#              {"title": "Slide 5", "type": "2"}
#          ]
#          },
#         {"title": "module 2",
#          "number": "2",
#          "slides": [
#              {"title": "Slide 1", "type": "2"},
#              {"title": "Slide 2", "type": "2"},
#              {"title": "Slide 3", "type": "2"},
#              {"title": "Slide 4", "type": "2"},
#              {"title": "Slide 5", "type": "2"}
#          ]
#          },
#         {"title": "module 3",
#          "number": "3",
#          "slides": [
#              {"title": "Slide 1", "type": "2"},
#              {"title": "Slide 2", "type": "2"},
#              {"title": "Slide 3", "type": "2"},
#              {"title": "Slide 4", "type": "2"},
#              {"title": "Slide 5", "type": "2"}
#          ]
#          }
#
#     ]
}


# iterate through modules
# for module in data['modules']:
#     # print module['title']
#
#     # create title slide for module
#     title_slide = prs.slides.add_slide(layout_title)
#     title = title_slide.shapes.title
#
#     # add title to slide
#     title.text = module['title']
#
#     # create section headers for module
#     for slide in module['slides']:
#         # print slide
#
#         # create slides for sections
#         section = prs.slides.add_slide(layout_subheader)
#         section_title = section.shapes.title
#         section_title.text = slide['title']
#
#     # save module file
#     prs.save(courseSlug+'-m'+module['number']+'.pptx')
#
#     print "created slides for module " + module['number']


prs = Presentation('ps-template.pptx')
prev_module = 0;
last_item = False

for i, slide in enumerate(data['slides']):

    print "working on item "+str(i)

    if i == len(data['slides'])-1:
        last_item = True

    cur_module = slide['module_num']

    # new module
    # save current preso, create new, add slide
    if last_item:

        # add slide
        section = prs.slides.add_slide(prs.slide_layouts[slide['slide_layout']])
        section_title = section.shapes.title
        section_title.text = slide['slide_title']

        prs.save(courseSlug + '-m' + str(cur_module) + '.pptx')
        print "created slides for module " + str(cur_module)

    elif cur_module == prev_module or prev_module == 0:

        # add slide
        section = prs.slides.add_slide(prs.slide_layouts[slide['slide_layout']])
        section_title = section.shapes.title
        section_title.text = slide['slide_title']

    else:
        # save current preso
        prs.save(courseSlug + '-m' + str(prev_module) + '.pptx')

        # print msg
        print "created slides for module " + str(prev_module)

        # create new preso
        prs = Presentation('ps-template.pptx')

        # add slide
        section = prs.slides.add_slide(prs.slide_layouts[slide['slide_layout']])
        section_title = section.shapes.title
        section_title.text = slide['slide_title']

    #set prev_module for iteration
    prev_module = cur_module
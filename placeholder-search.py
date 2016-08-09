from pptx import Presentation

prs = Presentation('ps-template.pptx')

slide = prs.slides.add_slide(prs.slide_layouts[0])

for shape in slide.shapes:
    if shape.is_placeholder:
        phf = shape.placeholder_format

        if phf.type != 'PICTURE (18)':
            txt = 'pic'
        else:
            txt = shape.text

        print('%d, %s, %s' % (phf.idx, phf.type, txt))
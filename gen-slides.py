from pptx import Presentation
import os
import config

# Placeholders in slide 0
# 12, BODY (2)
# 14, BODY (2)
# 10, BODY (2)
# 15, PICTURE (18)
# 13, BODY (2)
# 0, TITLE (1)

def checkOrCreateOutputDir():
    outputDir = config.outputDir
    if not os.path.exists(outputDir):
        os.makedirs(outputDir)

def addAuthorImage(section):
    headerImg = config.headerImg
    placeholder = section.placeholders[15]
    picture = placeholder.insert_picture(headerImg)
    section.placeholders[12].text = config.authorFooter
    section.placeholders[14].text = config.authorTitle
    section.placeholders[10].text = config.authorName


def buildSlides():
    courseName = config.courseName
    courseSlug = config.courseSlug
    outputDir = config.outputDir
    headerImg = config.headerImg
    authorName = config.authorName
    authorTitle = config.authorTitle
    authorFooter = config.authorFooter
    data = config.data

    prs = Presentation(config.templatePath)

    prev_module = 0;
    last_item = False

    for i, slide in enumerate(data['slides']):

        print "working on item "+str(i)

        if i == len(data['slides'])-1:
            last_item = True

        cur_module = slide['module_num']

        # new module
        # save current preso, create new, add slide
        if last_item:

            # add slide
            section = prs.slides.add_slide(prs.slide_layouts[slide['slide_layout']])
            section_title = section.shapes.title
            section_title.text = slide['slide_title']

            prs.save(outputDir + courseSlug + '-m' + str(cur_module) + '.pptx')
            print "created slides for module " + str(cur_module)

        elif cur_module == prev_module or prev_module == 0:

            # add slide
            section = prs.slides.add_slide(prs.slide_layouts[slide['slide_layout']])
            section_title = section.shapes.title
            section_title.text = slide['slide_title']

            # add author image
            if slide['slide_layout'] == 0:
                addAuthorImage(section)


        else:
            # save current preso
            prs.save(outputDir + courseSlug + '-m' + str(prev_module) + '.pptx')

            # print msg
            print "created slides for module " + str(prev_module)

            # create new preso
            prs = Presentation(config.templatePath)

            # add slide
            section = prs.slides.add_slide(prs.slide_layouts[slide['slide_layout']])
            section_title = section.shapes.title
            section_title.text = slide['slide_title']

            # add author image
            if slide['slide_layout'] == 0:
                addAuthorImage(section)


        #set prev_module for iteration
        prev_module = cur_module

# Execution Stack
checkOrCreateOutputDir()
buildSlides()
